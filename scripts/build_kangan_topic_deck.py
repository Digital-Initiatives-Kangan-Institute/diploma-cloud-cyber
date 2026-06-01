#!/usr/bin/env python3
"""Experiment: a Kangan-Institute-branded teaching deck (Topic 2 — The case for change).

The teaching deck is a REAL Kangan asset, so it wears Kangan/BKI branding — NOT the
in-world YAT case-study brand (YAT branding stays inside the scenario documents).

Brand extracted from the live kangan.edu.au `.kangan-theme` CSS variables:
  primary gold #edab0c (dark #d68e10, bright #fbb900) · charcoal #2a2929 on white ·
  greys #484848 / #7a7a7a · category accents magenta #92268f, sky #27b5ce, green #205f61,
  navy #004488 · type = Roboto (Light/Medium/Bold).

Usage:  python scripts/build_kangan_topic_deck.py [output.pptx]
Default: ../diploma-cloud-cyber/S1-CL1-Cloud-Design-Build/delivery/topic_02/Topic_02_Kangan_experiment.pptx
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- brand ----
GOLD     = "EDAB0C"
GOLD_DK  = "D68E10"
GOLD_BR  = "FBB900"
CHAR     = "2A2929"   # charcoal
INK      = "000000"
GREY1    = "484848"
GREY2    = "7A7A7A"
WHITE    = "FFFFFF"
BGLIGHT  = "F9F9F9"
BORDER   = "CCCCCC"
# category accents
MAGENTA  = "92268F"
SKY      = "27B5CE"
GREEN    = "205F61"
NAVY     = "004488"

FONT_BOLD = "Roboto"      # rendered Bold
FONT_MED  = "Roboto"      # rendered Medium/Regular
FONT_LT   = "Roboto Light"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _rgb(h):
    return RGBColor.from_string(h)


def _bg(slide, hexcol):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = _rgb(hexcol)


def _rect(slide, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(line); sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _box(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tb, tf


def _run(p, text, size, color, bold=False, italic=False, font=FONT_MED):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = _rgb(color)
    return r


def _para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p


def _footer(slide, pageno, accent=GOLD):
    # thin accent rule + wordmark + page number
    _rect(slide, Inches(0), Inches(7.18), Inches(13.333), Pt(3), fill=accent)
    tb, tf = _box(slide, Inches(0.55), Inches(7.0), Inches(6), Inches(0.35))
    p = _para(tf, True)
    _run(p, "Kangan ", 11, CHAR, bold=True, font=FONT_BOLD)
    _run(p, "Institute", 11, GOLD_DK, bold=True, font=FONT_BOLD)
    tb2, tf2 = _box(slide, Inches(11.5), Inches(7.0), Inches(1.3), Inches(0.35))
    p2 = _para(tf2, True); p2.alignment = PP_ALIGN.RIGHT
    _run(p2, str(pageno), 11, GREY2, font=FONT_MED)


def _bullets(tf, items, base_size=18):
    """items: list of (level, text, opts). opts: dict(bold,color,italic,marker)."""
    for i, it in enumerate(items):
        lvl, text = it[0], it[1]
        opts = it[2] if len(it) > 2 else {}
        p = _para(tf, first=(i == 0))
        p.space_after = Pt(10 if lvl == 0 else 5)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        size = base_size - (3 if lvl >= 1 else 0) - (3 if lvl >= 2 else 0)
        marker = opts.get("marker")
        if marker is None:
            marker = "" if lvl == 0 else ("–  " if lvl == 1 else "·  ")
        indent = Inches(0.0 if lvl == 0 else (0.4 if lvl == 1 else 0.8))
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(int(indent)))
        pPr.set('indent', '0')
        if lvl == 0:
            # leading gold square marker for top-level points
            _run(p, "■  ", size, opts.get("mark_color", GOLD), bold=True, font=FONT_BOLD)
        elif marker:
            _run(p, marker, size, GREY2, font=FONT_MED)
        _run(p, text, size, opts.get("color", CHAR),
             bold=opts.get("bold", False), italic=opts.get("italic", False),
             font=FONT_BOLD if opts.get("bold") else FONT_MED)


# ---------- slide builders ----------
def title_slide(prs, topic_no, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, CHAR)
    _rect(s, Inches(0), Inches(0), Inches(0.45), Inches(7.5), fill=GOLD)
    tb, tf = _box(s, Inches(1.0), Inches(2.2), Inches(11.5), Inches(2.6))
    p = _para(tf, True)
    _run(p, f"TOPIC {topic_no}", 20, GOLD_BR, bold=True, font=FONT_BOLD)
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    _run(p2, title, 46, WHITE, bold=True, font=FONT_BOLD)
    p3 = tf.add_paragraph(); p3.space_before = Pt(12)
    _run(p3, subtitle, 22, "D1D1D1", italic=True, font=FONT_LT)
    tb2, tf2 = _box(s, Inches(1.0), Inches(6.5), Inches(8), Inches(0.5))
    pp = _para(tf2, True)
    _run(pp, "Kangan ", 16, WHITE, bold=True, font=FONT_BOLD)
    _run(pp, "Institute", 16, GOLD, bold=True, font=FONT_BOLD)
    return s


def divider_slide(prs, number, title, kicker, accent):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _rect(s, Inches(0), Inches(0), Inches(4.3), Inches(7.5), fill=accent)
    # big number
    tb, tf = _box(s, Inches(0.4), Inches(2.3), Inches(3.6), Inches(2.4))
    p = _para(tf, True)
    _run(p, number, 130, WHITE, bold=True, font=FONT_BOLD)
    # title block on white
    tb2, tf2 = _box(s, Inches(4.9), Inches(2.7), Inches(7.8), Inches(2.4))
    p1 = _para(tf2, True)
    _run(p1, "SECTION", 16, accent, bold=True, font=FONT_BOLD)
    p2 = tf2.add_paragraph(); p2.space_before = Pt(6)
    _run(p2, title, 38, CHAR, bold=True, font=FONT_BOLD)
    if kicker:
        p3 = tf2.add_paragraph(); p3.space_before = Pt(10)
        _run(p3, kicker, 20, GREY1, italic=True, font=FONT_LT)
    _footer(s, "", accent=accent)
    return s


def content_slide(prs, pageno, title, kicker, bullets, accent=GOLD, base=18):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    # title
    tb, tf = _box(s, Inches(0.7), Inches(0.55), Inches(11.9), Inches(1.1))
    p = _para(tf, True)
    _run(p, title, 30, CHAR, bold=True, font=FONT_BOLD)
    if kicker:
        pk = tf.add_paragraph(); pk.space_before = Pt(2)
        _run(pk, kicker, 18, GOLD_DK, italic=True, font=FONT_LT)
    # gold rule under title
    _rect(s, Inches(0.72), Inches(1.62), Inches(1.6), Pt(4), fill=accent)
    # body
    tb2, tf2 = _box(s, Inches(0.72), Inches(1.95), Inches(11.9), Inches(4.8))
    _bullets(tf2, bullets, base_size=base)
    _footer(s, pageno, accent=accent)
    return s


def activity_slide(prs, pageno, title, bullets, timer, accent=GOLD):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    # gold header band
    _rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.35), fill=accent)
    # ACTIVITY pill
    pill = _rect(s, Inches(0.7), Inches(0.42), Inches(1.85), Inches(0.5),
                 fill=CHAR, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ptf = pill.text_frame; ptf.word_wrap = False
    pa = ptf.paragraphs[0]; pa.alignment = PP_ALIGN.CENTER
    _run(pa, "ACTIVITY", 14, GOLD_BR, bold=True, font=FONT_BOLD)
    tb, tf = _box(s, Inches(2.8), Inches(0.34), Inches(9.8), Inches(0.75), anchor=MSO_ANCHOR.MIDDLE)
    p = _para(tf, True)
    _run(p, title, 26, WHITE, bold=True, font=FONT_BOLD)
    # body
    tb2, tf2 = _box(s, Inches(0.72), Inches(1.7), Inches(11.9), Inches(4.7))
    _bullets(tf2, bullets, base_size=17)
    # timer strip
    _rect(s, Inches(0.7), Inches(6.45), Inches(4.4), Inches(0.5),
          fill=BGLIGHT, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tbt, tft = _box(s, Inches(0.9), Inches(6.5), Inches(4.0), Inches(0.4), anchor=MSO_ANCHOR.MIDDLE)
    pt = _para(tft, True)
    _run(pt, "⏱  " + timer, 15, CHAR, bold=True, font=FONT_BOLD)
    _footer(s, pageno, accent=accent)
    return s


def takeaways_slide(prs, pageno, section_label, points, accent=GOLD):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BGLIGHT)
    # header
    tb, tf = _box(s, Inches(0.7), Inches(0.6), Inches(11.9), Inches(1.2))
    p = _para(tf, True)
    _run(p, "Key takeaways", 32, CHAR, bold=True, font=FONT_BOLD)
    pk = tf.add_paragraph(); pk.space_before = Pt(2)
    _run(pk, section_label, 18, accent, bold=True, font=FONT_BOLD)
    _rect(s, Inches(0.72), Inches(1.75), Inches(1.6), Pt(4), fill=accent)
    # points as cards
    top = Inches(2.15)
    gap = Inches(0.12)
    card_h = Inches(0.92)
    for i, txt in enumerate(points):
        y = Emu(int(top) + i * (int(card_h) + int(gap)))
        _rect(s, Inches(0.7), y, Inches(11.95), card_h, fill=WHITE, line=BORDER)
        _rect(s, Inches(0.7), y, Inches(0.12), card_h, fill=accent)
        tbc, tfc = _box(s, Inches(1.05), y, Inches(11.4), card_h, anchor=MSO_ANCHOR.MIDDLE)
        pc = _para(tfc, True); pc.line_spacing = 1.0
        _run(pc, txt, 17, CHAR, font=FONT_MED)
    _footer(s, pageno, accent=accent)
    return s


def table_slide(prs, pageno, title, kicker, headers, rows, accent=GOLD,
                col_widths=None, note=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    tb, tf = _box(s, Inches(0.7), Inches(0.55), Inches(11.9), Inches(1.1))
    p = _para(tf, True)
    _run(p, title, 30, CHAR, bold=True, font=FONT_BOLD)
    if kicker:
        pk = tf.add_paragraph(); pk.space_before = Pt(2)
        _run(pk, kicker, 18, GOLD_DK, italic=True, font=FONT_LT)
    _rect(s, Inches(0.72), Inches(1.62), Inches(1.6), Pt(4), fill=accent)
    nrows, ncols = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(nrows, ncols, Inches(0.72), Inches(2.0),
                            Inches(11.9), Inches(0.5 * nrows)).table
    gt.first_row = False; gt.horz_banding = False
    if col_widths:
        for c, w in enumerate(col_widths):
            gt.columns[c].width = Inches(w)
    for c, htext in enumerate(headers):
        cell = gt.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(accent)
        cell.margin_left = cell.margin_right = Pt(6)
        para = cell.text_frame.paragraphs[0]
        _run(para, htext, 14, WHITE, bold=True, font=FONT_BOLD)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(WHITE if r % 2 else BGLIGHT)
            cell.margin_left = cell.margin_right = Pt(6)
            cell.margin_top = cell.margin_bottom = Pt(4)
            para = cell.text_frame.paragraphs[0]
            # first column of a field/value table is a label -> bold
            bold = (c == 0 and len(headers) == 2 and headers[0] in ("Field",))
            _run(para, val, 13, CHAR, bold=bold, font=FONT_BOLD if bold else FONT_MED)
    if note:
        tbn, tfn = _box(s, Inches(0.72), Inches(6.5), Inches(11.9), Inches(0.5))
        pn = _para(tfn, True)
        _run(pn, note, 15, GREY1, italic=True, font=FONT_LT)
    _footer(s, pageno, accent=accent)
    return s


# ---------- content (Topic 2) ----------
def build(path):
    prs = Presentation()
    prs.slide_width = EMU_W; prs.slide_height = EMU_H
    n = 0
    def pg():
        nonlocal n; n += 1; return n

    # Title
    title_slide(prs, "02", "The case for change", "Prove the problem before you propose the fix")

    # Opener framing
    content_slide(prs, pg(), "The consultant's first real job", "diagnose before you prescribe", [
        (0, "You can speak cloud now — now you do the diagnosis."),
        (0, "Find out where the organisation is, where it needs to be, and the gap between."),
        (0, "Three moves: align to strategy → describe the current state → expose the gaps.",
            {"bold": True, "color": GOLD_DK}),
        (0, "These three become the opening sections of the business case you'll build."),
    ])
    content_slide(prs, pg(), "A case for change is evidence, not opinion", "", [
        (0, "It answers three questions, in order:"),
        (1, "Where does the organisation want to go?   — strategy"),
        (1, "Where is it now?   — current state"),
        (1, "What's missing to get there?   — gaps"),
        (0, "You're not choosing a solution yet — you're proving the problem is real and worth solving."),
        (0, "Solutions come next. Diagnose before you prescribe.", {"bold": True, "color": MAGENTA, "mark_color": MAGENTA}),
    ])

    # ---- Section 1 ----
    divider_slide(prs, "01", "Strategic alignment", "tie the initiative to the organisation's goals", MAGENTA)
    content_slide(prs, pg(), "Start with strategy", "or the board won't fund it", [
        (0, "A board funds change that moves the organisation toward its own stated goals."),
        (0, "Lead by showing the initiative serves the strategy — not technology for its own sake."),
        (0, "Strategic alignment = “here's what they said they want; here's how this delivers it.”"),
        (0, "It's the “why this, why now” that earns everything that follows.", {"bold": True, "color": GOLD_DK}),
    ], accent=MAGENTA)
    content_slide(prs, pg(), "Read a strategic plan in three layers", "", [
        (0, "Business objectives — where the whole organisation is going",
            {"bold": True}),
        (1, "e.g. grow students 15%/yr, expand nationally"),
        (0, "ICT goals / objectives — how ICT supports that", {"bold": True}),
        (1, "e.g. reduce in-house dependency; 99.9% for critical systems"),
        (0, "Initiatives — the planned actions", {"bold": True}),
        (1, "e.g. move suitable on-site systems to the cloud"),
        (0, "Trace your initiative up these layers — pull only what's material.",
            {"bold": True, "color": MAGENTA, "mark_color": MAGENTA}),
    ], accent=MAGENTA)
    content_slide(prs, pg(), "Add the outside view", "industry context", [
        (0, "A plan doesn't exist in a vacuum — compare it to where the industry is heading."),
        (0, "Trends that matter: cloud adoption · managed services · OPEX over CAPEX · resilience."),
        (0, "Name both directions:"),
        (1, "Alignment — the plan matches industry direction (strengthens the case)"),
        (1, "Divergence — the plan lags or differs (a risk, or an opportunity to flag)"),
    ], accent=MAGENTA)
    content_slide(prs, pg(), "What good looks like", "", [
        (0, "Every claim cited to the plan (e.g. “ICT Strategic Plan — ICT Goals”)."),
        (0, "Material items only — not the whole plan recited."),
        (0, "Traces the initiative up to a real organisational goal."),
        (0, "Alignment and divergence both named — balanced, not cheerleading.",
            {"bold": True, "color": MAGENTA, "mark_color": MAGENTA}),
    ], accent=MAGENTA)
    activity_slide(prs, pg(), "Write a Strategic Alignment section", [
        (0, "In your working copy of the Business Case, write the Strategic Alignment section for moving the Accounting System (Ledgerline).", {"bold": True}),
        (0, "Remember to:"),
        (1, "Read YAT's ICT Strategic Plan."),
        (1, "Which ICT goals/objectives does this migration serve? Trace up to a business objective — cite them."),
        (1, "Add one or two points of industry context — where's the sector heading?"),
        (1, "Note alignment and divergence."),
        (0, "~½ page."),
        (0, "⚠  Argue alignment through the ICT goals, and name the weaker links honestly. That is the analysis.",
            {"color": MAGENTA, "mark_color": MAGENTA}),
    ], "~20 min, then we discuss", accent=MAGENTA)
    takeaways_slide(prs, pg(), "Section 1 · Strategic alignment", [
        "Open with strategy: tie the initiative to the organisation's own stated goals.",
        "Three layers — business objectives / ICT goals / initiatives; pull only what's material.",
        "Add the outside view; name alignment and divergence.",
        "Cite everything. Strategy is the “why this, why now”.",
    ], accent=MAGENTA)

    # ---- Section 2 ----
    divider_slide(prs, "02", "Current-state synthesis", "synthesis, not transcription", SKY)
    content_slide(prs, pg(), "Current state is synthesis, not transcription", "", [
        (0, "Summarise the current environment in your own words, focused on what's material to this decision."),
        (0, "The board needs the picture that explains why change is needed — not a copy of the IT docs."),
        (0, "Distil it down; don't transcribe it.", {"bold": True, "color": GOLD_DK}),
    ], accent=SKY)
    content_slide(prs, pg(), "The relevance filter", "keep vs drop", [
        (0, "KEEP (if it bears on the decision):", {"bold": True}),
        (1, "platform & stack · age / condition / capacity · availability today · dependencies & integrations · constraints · pain points"),
        (0, "DROP:", {"bold": True}),
        (1, "detail that doesn't move the decision (printer counts, unrelated systems)"),
        (0, "A good current state quietly surfaces the limitations that motivate the change.",
            {"bold": True, "color": SKY, "mark_color": SKY}),
    ], accent=SKY)
    content_slide(prs, pg(), "How to distil", "", [
        (0, "Read the source docs — environment overview, server/app specs, consultation notes.", {"marker": "1.  "}),
        (0, "For each fact ask: “does this affect the renew-vs-migrate decision?” If no, cut it.", {"marker": "2.  "}),
        (0, "Re-state what's left in plain language, grouped: platform · workload · dependencies · condition/risk.", {"marker": "3.  "}),
        (0, "Aim ~½ page. No copy-paste.", {"bold": True, "color": GOLD_DK}),
    ], accent=SKY)
    activity_slide(prs, pg(), "Add the Current State section", [
        (0, "In your working copy of the Business Case, add the Current State section for the Accounting System (Ledgerline).", {"bold": True}),
        (0, "Remember to:"),
        (1, "Use the Accounting application spec, server specs, operational costing, the ICT Environment Overview, and the consultation notes."),
        (1, "Cover: platform & stack · workload (incl. month-end / EOFY peaks) · integrations (AD, O365, LMS fee-status, payroll, banking) · condition & constraints."),
        (0, "~½ page."),
        (0, "⚠  Material only — synthesise in your own words, no verbatim copying.",
            {"color": SKY, "mark_color": SKY}),
    ], "~20 min, then we discuss", accent=SKY)
    takeaways_slide(prs, pg(), "Section 2 · Current state", [
        "Own words, material facts only — synthesis beats transcription.",
        "Filter every fact against “does it affect the decision?”",
        "Surface the limitations that set up the gap analysis.",
    ], accent=SKY)

    # ---- Section 3 ----
    divider_slide(prs, "03", "Gap analysis", "bridges “want” and “have”", GREEN)
    content_slide(prs, pg(), "Gap analysis bridges “want” and “have”", "", [
        (0, "It bridges where we want to be (strategy + requirements) and where we are (current state)."),
        (0, "It makes the problem concrete and measurable."),
        (0, "It hands the next stage a list of changes to evaluate."),
        (0, "No gap → no case.", {"bold": True, "color": GREEN, "mark_color": GREEN}),
    ], accent=GREEN)
    table_slide(prs, pg(), "The gap table", "one row per objective",
                ["Objective", "Current", "Desired", "Gap", "Opportunity", "Proposed change"],
                [["", "", "", "", "", ""]],
                accent=GREEN,
                col_widths=[2.2, 2.0, 1.9, 1.9, 1.9, 2.0],
                note="Objective/Desired ← Strategic Alignment + requirements · Current ← your Current State · Proposed change = the seed of the options you'll weigh next.")
    table_slide(prs, pg(), "Anatomy of a good row", "one objective, traced end to end",
                ["Field", "Example"],
                [["Objective", "“reduce dependency on in-house server infrastructure”"],
                 ["Current", "single ageing on-prem server, owned and maintained in-house"],
                 ["Desired", "no in-house server to own / patch / replace"],
                 ["Gap", "full reliance on end-of-life hardware the organisation must run itself"],
                 ["Opportunity", "move to a managed / cloud platform"],
                 ["Proposed change", "evaluate migrating the workload to the cloud (→ next topic)"]],
                accent=GREEN,
                col_widths=[2.6, 9.3],
                note="Traceable both ways — objective came from the plan, current from your Current State.")
    activity_slide(prs, pg(), "Build the Gap Analysis", [
        (0, "In your working copy of the Business Case, add the Gap Analysis for the Accounting System (Ledgerline) — a table of at least 3 rows.", {"bold": True}),
        (0, "Remember to draw objectives from your Strategic Alignment work + the migration requirements, e.g.:"),
        (1, "reduce in-house infrastructure dependency (ageing server)"),
        (1, "business-hours availability ≥ 99.5% (RPO ≤ 1h / RTO ≤ 1 business day)"),
        (1, "keep financial data onshore + 7-year retention"),
        (1, "size for month-end / EOFY peaks without year-round over-provisioning"),
        (0, "Fill every column: objective → current → desired → gap → opportunity → proposed change."),
        (0, "⚠  Each row must trace back to a real objective and your current-state facts.",
            {"color": GREEN, "mark_color": GREEN}),
    ], "~25 min, then we discuss", accent=GREEN)
    takeaways_slide(prs, pg(), "Section 3 · Gap analysis", [
        "Gap analysis = desired (strategy + requirements) vs current (your synthesis).",
        "One row per objective; fill all six columns; keep it traceable.",
        "The “proposed change” column feeds straight into the options analysis next.",
    ], accent=GREEN)

    # ---- Close ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, CHAR)
    _rect(s, Inches(0), Inches(0), Inches(0.45), Inches(7.5), fill=GOLD)
    tb, tf = _box(s, Inches(1.0), Inches(1.5), Inches(11.5), Inches(4.5))
    p = _para(tf, True)
    _run(p, "From diagnosis to evaluation", 40, WHITE, bold=True, font=FONT_BOLD)
    for txt in [
        "You've built the case for change: strategy → current state → gaps.",
        "That's the first half of a business case — it proves the problem is real and worth solving.",
        "Your “proposed changes” aren't decisions yet — they're candidates.",
        "Next: weigh the options, cost them, assess the risk — turning the case for change into a recommendation.",
        "You practised on the Accounting System; the same moves carry to whatever engagement lands on your desk.",
    ]:
        pp = tf.add_paragraph(); pp.space_before = Pt(12); pp.line_spacing = 1.05
        _run(pp, txt, 19, "E8E8E8", font=FONT_LT)

    Path(path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    print(f"Wrote {path} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    default = "S1-CL1-Cloud-Design-Build/delivery/topic_02/Topic_02_Kangan_experiment.pptx"
    out = sys.argv[1] if len(sys.argv) > 1 else default
    build(out)
